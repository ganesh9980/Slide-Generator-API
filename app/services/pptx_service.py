from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
import os
from ..config import Config
from ..models.presentation import Presentation as PresModel, Slide


class PPTXService:
    def __init__(self):
        self.template_path = os.path.join(Config.TEMPLATE_FOLDER, 'default.pptx')
        self.output_folder = Config.OUTPUT_FOLDER

        # Create directories if missing
        os.makedirs(Config.TEMPLATE_FOLDER, exist_ok=True)
        os.makedirs(self.output_folder, exist_ok=True)

        # Create blank template if none exists
        if not os.path.exists(self.template_path):
            Presentation().save(self.template_path)

    def generate_pptx(self, presentation: PresModel) -> str:
        """
        Generate PowerPoint file from presentation model with enhanced content
        """
        prs = Presentation(self.template_path)

        # Apply theme
        primary_color = self._hex_to_rgb(presentation.theme.get('primary_color', '#2A5CAA'))
        secondary_color = self._hex_to_rgb(presentation.theme.get('secondary_color', '#5A8F29'))
        font_name = presentation.theme.get('font', 'Calibri')

        # Process each slide
        for slide_data in presentation.slides:
            self._add_enhanced_slide(prs, slide_data, primary_color, secondary_color, font_name)

        # Save file
        file_path = os.path.join(self.output_folder, f"{presentation.id}.pptx")
        prs.save(file_path)
        return file_path

    def _add_enhanced_slide(self, prs, slide: Slide, primary_color, secondary_color, font_name):
        """Create slide from Slide object"""
        if slide.layout == 'title':
            layout = prs.slide_layouts[0]  # Title slide
        else:
            layout = prs.slide_layouts[1]  # Content slide

        slide_obj = prs.slides.add_slide(layout)

        # Set title
        title = slide_obj.shapes.title
        title.text = slide.title
        title.text_frame.paragraphs[0].font.color.rgb = primary_color
        title.text_frame.paragraphs[0].font.name = font_name

        # Add content points
        if hasattr(slide, 'points') and slide.points and len(slide_obj.shapes.placeholders) > 1:
            content = slide_obj.shapes.placeholders[1].text_frame
            for point in slide.points:
                p = content.add_paragraph()
                p.text = point
                p.font.name = font_name
                p.font.size = Pt(18)
                p.level = 0

        # Add image placeholder
        if slide.image_suggestion:
            self._add_image_placeholder(slide_obj, slide.image_suggestion, secondary_color)

        # Add citation footer
        if slide.citation:
            self._add_citation_footer(slide_obj, slide.citation, font_name)

    def _add_image_placeholder(self, slide, image_query, color):
        """Add image placeholder with search suggestion"""
        left = Inches(5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(2.5)

        # Create placeholder rectangle
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = color
        placeholder.line.width = Pt(1.5)

        # Add search suggestion text
        textbox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Suggested image: {image_query} (CC license)"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 100)

    def _add_citation_footer(self, slide, citation, font_name):
        """Add citation at bottom of slide"""
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Source: {citation}"
        p.font.size = Pt(10)
        p.font.name = font_name
        p.font.italic = True
        p.alignment = PP_ALIGN.RIGHT

    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(*tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4)))